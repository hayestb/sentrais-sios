#!/usr/bin/env python3
"""
Sentrais SIOS — Legacy Document Ingestion Pipeline
Organizes, extracts, chunks, embeds, and loads documents into Neon (pgvector).

Usage:
  python3 scripts/ingest-docs.py --source "/path/to/files (13)" [--dry-run] [--reset]

Requires env vars (copy from .env.local):
  DATABASE_URL_UNPOOLED   — direct Neon connection (no pooler)
  OPENAI_API_KEY          — for text-embedding-3-small (1536 dims)
"""

import argparse
import hashlib
import json
import os
import re
import shutil
import sys
from pathlib import Path

# Auto-load .env.local from the project root (two levels up from scripts/)
_env_path = Path(__file__).resolve().parent.parent / ".env.local"
if _env_path.exists():
    from dotenv import load_dotenv
    load_dotenv(_env_path, override=True)
import time
from typing import Optional

# ── Dependencies ──────────────────────────────────────────────────────────────
try:
    import pdfplumber
    import docx
    import pptx
    import tiktoken
    import psycopg2
    import psycopg2.extras
    import voyageai
except ImportError as e:
    print(f"Missing dependency: {e}")
    print("Run: pip3 install pdfplumber python-docx python-pptx tiktoken psycopg2-binary voyageai")
    sys.exit(1)

# ── Config ────────────────────────────────────────────────────────────────────

CHUNK_TOKENS = 512       # target tokens per chunk
CHUNK_OVERLAP = 64       # overlap between adjacent chunks
EMBED_MODEL   = "voyage-3-lite"
EMBED_DIMS    = 512
EMBED_BATCH   = 128      # voyage-3-lite supports up to 128 texts per batch

CATEGORIES: dict[str, list[str]] = {
    "sentrais-core": [
        "sentrais", "brand", "capabilities", "federal", "blueprint", "cpo",
        "enterprise", "architecture", "handshake", "gartner", "chimera",
        "novatelabs", "novate", "resilience orchestration", "business plan",
        "investor", "pitch", "one pager", "onepager",
    ],
    "nfl": [
        "nfl", "league", "football", "nflit", "hub", "nin",
        "innovation framework", "innovation hub", "arb", "game day",
        "internship", "scouting", "tentpole",
    ],
    "evergame": [
        "evergame", "evergame360", "evergame 360", "evergame master",
        "eg360", "eg 360", "rapid build",
    ],
    "spectra-civigrid": [
        "spectra", "civigrid", "evenu", "evenu360", "epicenter",
        "esports", "e-sports", "gaming",
    ],
    "legal-contracts": [
        "msa", "term sheet", "agreement", "offer letter", "attorney",
        "draft", "contract", "consulting agreement", "employment",
        "sow", "statement of work", "invoice", "ein",
    ],
    "operations": [
        "playbook", "raci", "operations", "sear", "leors", "sipe",
        "go-to-market", "framework", "automation", "orchestration",
        "command", "unified ops", "deployment", "transition", "cadence",
        "process map", "milestone", "execution", "gda", "poc",
    ],
    "personal": [
        "tye hayes", "speaker", "capability statement", "new orleans",
        "smartcity", "supercycle", "linkedin", "samson", "shephard",
    ],
}

# Files to skip entirely (security-sensitive or non-document)
SKIP_FILENAMES = {
    "github-recovery-codes - NOVATEHER.txt",
    "NordPass Recovery Code.pdf",
    ".DS_Store",
    "sentrais-pmo.bundle",
    # large/problematic files
    "PROFESSIONAL_SERVICES_AGREEMENT__Asia_Lewis.pdf",
    "PROFESSIONAL SERVICES AGREEMENT - Asia Lewis.pdf",
    "NFL Innovation Framework Final V5_5_1.pptx",
    "NFL_Innovation_Framework_Final_V5_5_1.pptx",
    "NFL Innovation Hub",
}

MAX_FILE_SIZE_MB = 20  # skip files larger than this

SKIP_EXTENSIONS = {
    ".zip", ".png", ".jpg", ".jpeg", ".gif", ".svg", ".bundle",
    ".html", ".csv", ".xlsx", ".xls", ".numbers", ".pages",
    ".key", ".app", ".dmg", ".pkg", ".textClipping",
    ".json", ".ts", ".tsx", ".js", ".jsx", ".py", ".sh",
    ".env", ".lock", ".log", ".map", ".tsbuildinfo",
}

# Only ingest these document extensions
INGEST_EXTENSIONS = {".pdf", ".docx", ".pptx", ".txt", ".md"}  # .doc unsupported by python-docx

# Folder names to skip entirely when recursing
SKIP_DIRS = {
    "node_modules", ".git", ".next", "__pycache__", ".cache",
    "dist", "build", ".turbo", "coverage", "venv", ".venv",
}

# ── Helpers ───────────────────────────────────────────────────────────────────

enc = tiktoken.get_encoding("cl100k_base")


def classify(filename: str) -> str:
    lower = filename.lower()
    for category, keywords in CATEGORIES.items():
        for kw in keywords:
            if kw.lower() in lower:
                return category
    return "other"


def sanitize_filename(name: str) -> str:
    name = name.replace("%20", " ")
    name = re.sub(r"\.pdf\.pdf$", ".pdf", name)
    name = re.sub(r"[^\w\s\-_\.]", "_", name)
    name = re.sub(r"\s+", "_", name.strip())
    return name


def extract_text_pdf(path: Path) -> tuple[str, int]:
    import signal

    def _timeout(signum, frame):
        raise TimeoutError("PDF extraction timed out")

    pages = []
    try:
        signal.signal(signal.SIGALRM, _timeout)
        signal.alarm(30)  # 30-second limit per PDF
        try:
            with pdfplumber.open(path) as pdf:
                for page in pdf.pages:
                    text = page.extract_text() or ""
                    pages.append(text)
        finally:
            signal.alarm(0)
        return "\n\n".join(pages), len(pages)
    except TimeoutError:
        print(f"    ⚠ PDF timed out ({path.name}) — skipping")
        return "", 0
    except Exception as e:
        print(f"    ⚠ PDF extract failed ({path.name}): {e}")
        return "", 0


def extract_text_docx(path: Path) -> tuple[str, int]:
    try:
        doc = docx.Document(str(path))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n\n".join(paragraphs), 1
    except Exception as e:
        print(f"    ⚠ DOCX extract failed ({path.name}): {e}")
        return "", 0


def extract_text_pptx(path: Path) -> tuple[str, int]:
    try:
        prs = pptx.Presentation(str(path))
        slides = []
        for slide in prs.slides:
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text)
            slides.append("\n".join(texts))
        return "\n\n---\n\n".join(slides), len(prs.slides)
    except Exception as e:
        print(f"    ⚠ PPTX extract failed ({path.name}): {e}")
        return "", 0


def extract_text_txt(path: Path) -> tuple[str, int]:
    try:
        return path.read_text(errors="replace"), 1
    except Exception as e:
        print(f"    ⚠ TXT extract failed ({path.name}): {e}")
        return "", 0


def extract_text(path: Path) -> tuple[str, int]:
    ext = path.suffix.lower()
    if ext == ".pdf":
        return extract_text_pdf(path)
    elif ext == ".docx":
        return extract_text_docx(path)
    elif ext in (".pptx", ".ppt"):
        return extract_text_pptx(path)
    elif ext in (".txt", ".md"):
        return extract_text_txt(path)
    return "", 0


def chunk_text(text: str, doc_id: str, filename: str) -> list[dict]:
    tokens = enc.encode(text)
    chunks = []
    start = 0
    idx = 0
    while start < len(tokens):
        end = min(start + CHUNK_TOKENS, len(tokens))
        chunk_tokens = tokens[start:end]
        content = enc.decode(chunk_tokens)
        chunks.append({
            "chunk_index": idx,
            "content": content,
            "token_count": len(chunk_tokens),
            "metadata": {"source": filename, "chunk": idx},
        })
        idx += 1
        start += CHUNK_TOKENS - CHUNK_OVERLAP
    return chunks


def get_embeddings(client, texts: list[str]) -> list[list[float]]:
    """Embed in batches using Voyage AI."""
    all_embeddings = []
    for i in range(0, len(texts), EMBED_BATCH):
        batch = texts[i : i + EMBED_BATCH]
        resp = client.embed(batch, model=EMBED_MODEL, input_type="document")
        all_embeddings.extend(resp.embeddings)
        if i + EMBED_BATCH < len(texts):
            time.sleep(0.25)
    return all_embeddings


def file_hash(path: Path) -> str:
    # Use size + mtime as a fast proxy — avoids reading large files
    stat = path.stat()
    return f"{stat.st_size}:{stat.st_mtime}"


# ── Phase 1: Organize ─────────────────────────────────────────────────────────

def organize_files(source_dir: Path, dest_dir: Path, dry_run: bool) -> list[dict]:
    """Copy files into dest_dir/<category>/ with clean names. Returns manifest."""
    manifest = []
    dest_dir.mkdir(parents=True, exist_ok=True)

    # Collect all files recursively (skip dev dirs, non-doc extensions)
    all_files = []
    for f in sorted(source_dir.rglob("*")):
        if not f.is_file():
            continue
        # Skip files inside blacklisted directories anywhere in the path
        if any(part in SKIP_DIRS for part in f.parts):
            continue
        if f.name in SKIP_FILENAMES or sanitize_filename(f.name) in SKIP_FILENAMES:
            print(f"  ⏭  SKIP (blocked):   {f.name}")
            continue
        # Only ingest known document types
        if f.suffix.lower() not in INGEST_EXTENSIONS:
            continue
        size = f.stat().st_size
        if size == 0:
            print(f"  ⏭  SKIP (empty):     {f.name}")
            continue
        if size > MAX_FILE_SIZE_MB * 1024 * 1024:
            print(f"  ⏭  SKIP (too large): {f.name}")
            continue
        all_files.append(f)

    print(f"\n📂 Organizing {len(all_files)} files into {dest_dir}\n")

    for f in all_files:
        category = classify(f.name)
        clean_name = sanitize_filename(f.name)
        dest_cat = dest_dir / category
        dest_path = dest_cat / clean_name

        # Handle collisions
        if dest_path.exists() and file_hash(f) == file_hash(dest_path):
            print(f"  ✓ already exists: {category}/{clean_name}")
        elif dest_path.exists():
            stem, ext = os.path.splitext(clean_name)
            dest_path = dest_cat / f"{stem}_2{ext}"

        print(f"  → {category}/{clean_name}")
        if not dry_run:
            dest_cat.mkdir(parents=True, exist_ok=True)
            import signal as _sig
            def _timeout(s, f): raise TimeoutError()
            _sig.signal(_sig.SIGALRM, _timeout)
            _sig.alarm(15)
            try:
                shutil.copy2(f, dest_path)
            except TimeoutError:
                print(f"  ⚠ copy timed out — skipping {clean_name}")
                if dest_path.exists(): dest_path.unlink()
                continue
            finally:
                _sig.alarm(0)

        manifest.append({
            "original_path": str(f),
            "dest_path": str(dest_path),
            "filename": clean_name,
            "category": category,
            "mime_type": mime_for(f),
            "file_size_bytes": f.stat().st_size,
        })

    return manifest


def mime_for(path: Path) -> str:
    ext = path.suffix.lower()
    return {
        ".pdf": "application/pdf",
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ".txt": "text/plain",
        ".md": "text/markdown",
    }.get(ext, "application/octet-stream")


# ── Phase 2 & 3: Extract → Chunk → Embed → Load ──────────────────────────────

def ingest_to_db(manifest: list[dict], conn, voyage_client, dry_run: bool):
    cur = conn.cursor(cursor_factory=psycopg2.extras.RealDictCursor)

    for item in manifest:
        path = Path(item["dest_path"])
        print(f"\n📄 {item['category']}/{item['filename']}")

        # ── Extract ──
        text, page_count = extract_text(path)
        if not text.strip():
            print(f"    ⚠ no text extracted — skipping")
            continue
        print(f"    extracted {len(text):,} chars, {page_count} pages")

        # ── Insert document record ──
        if not dry_run:
            cur.execute(
                """
                INSERT INTO documents
                  (filename, original_path, category, mime_type, file_size_bytes,
                   status, extracted_text, page_count, metadata)
                VALUES (%s, %s, %s, %s, %s, 'processing', %s, %s, %s)
                ON CONFLICT DO NOTHING
                RETURNING id
                """,
                (
                    item["filename"],
                    item["original_path"],
                    item["category"],
                    item["mime_type"],
                    item["file_size_bytes"],
                    text,
                    page_count,
                    json.dumps({}),
                ),
            )
            row = cur.fetchone()
            if not row:
                # Already ingested — check if it has chunks
                cur.execute("SELECT id FROM documents WHERE filename = %s", (item["filename"],))
                existing = cur.fetchone()
                if existing:
                    cur.execute("SELECT COUNT(*) as n FROM document_chunks WHERE document_id = %s", (existing["id"],))
                    count = cur.fetchone()["n"]
                    if count > 0:
                        print(f"    ✓ already indexed ({count} chunks)")
                        continue
                    doc_id = str(existing["id"])
                else:
                    print(f"    ⚠ could not insert or find document record")
                    continue
            else:
                doc_id = str(row["id"])
            conn.commit()
        else:
            doc_id = "dry-run-id"

        # ── Chunk ──
        chunks = chunk_text(text, doc_id, item["filename"])
        print(f"    {len(chunks)} chunks")

        # ── Embed ──
        print(f"    embedding...", end=" ", flush=True)
        chunk_texts = [c["content"] for c in chunks]
        if not dry_run:
            embeddings = get_embeddings(voyage_client, chunk_texts)
        else:
            embeddings = [[0.0] * EMBED_DIMS for _ in chunks]
        print("done")

        # ── Insert chunks ──
        if not dry_run:
            records = [
                (
                    doc_id,
                    c["chunk_index"],
                    c["content"],
                    c["token_count"],
                    f"[{','.join(str(x) for x in emb)}]",
                    json.dumps(c["metadata"]),
                )
                for c, emb in zip(chunks, embeddings)
            ]
            psycopg2.extras.execute_values(
                cur,
                """
                INSERT INTO document_chunks
                  (document_id, chunk_index, content, token_count, embedding, metadata)
                VALUES %s
                """,
                records,
                template="(%s, %s, %s, %s, %s::vector, %s)",
            )
            cur.execute(
                "UPDATE documents SET status = 'indexed', updated_at = now() WHERE id = %s",
                (doc_id,),
            )
            conn.commit()
            print(f"    ✅ indexed")

    cur.close()


# ── CLI ───────────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(description="Sentrais legacy doc ingestion pipeline")
    parser.add_argument("--source", required=True, help="Source folder containing legacy docs")
    parser.add_argument("--dest", default=None, help="Destination folder (default: sentrais-sios/docs/knowledge-base)")
    parser.add_argument("--dry-run", action="store_true", help="Organize files and extract text, but skip DB writes")
    parser.add_argument("--reset", action="store_true", help="Delete all document_chunks and documents before re-ingesting")
    parser.add_argument("--category", default=None, help="Only ingest files matching this category")
    args = parser.parse_args()

    source_dir = Path(args.source).expanduser().resolve()
    if not source_dir.exists():
        print(f"Source not found: {source_dir}")
        sys.exit(1)

    script_dir = Path(__file__).resolve().parent
    dest_dir = Path(args.dest) if args.dest else script_dir.parent / "docs" / "knowledge-base"

    # ── Env checks ──
    db_url = os.environ.get("DATABASE_URL_UNPOOLED") or os.environ.get("DATABASE_URL")
    voyage_key = os.environ.get("VOYAGE_API_KEY")

    if not args.dry_run:
        if not db_url:
            print("Missing DATABASE_URL_UNPOOLED env var")
            sys.exit(1)
        if not voyage_key:
            print("Missing VOYAGE_API_KEY env var")
            sys.exit(1)

    print("=" * 60)
    print("SENTRAIS SIOS — Document Ingestion Pipeline")
    print("=" * 60)
    print(f"Source : {source_dir}")
    print(f"Dest   : {dest_dir}")
    print(f"Mode   : {'DRY RUN' if args.dry_run else 'LIVE'}")
    print()

    # ── Phase 1: Organize ──
    manifest = organize_files(source_dir, dest_dir, args.dry_run)

    if args.category:
        manifest = [m for m in manifest if m["category"] == args.category]
        print(f"\nFiltered to category '{args.category}': {len(manifest)} files")

    if args.dry_run:
        print(f"\n✅ Dry run complete. {len(manifest)} files would be ingested.")
        print("\nCategory breakdown:")
        from collections import Counter
        for cat, count in sorted(Counter(m["category"] for m in manifest).items()):
            print(f"  {cat}: {count}")
        return

    # ── Phase 2+3: Extract → Embed → Load ──
    voyage_client = voyageai.Client(api_key=voyage_key)
    conn = psycopg2.connect(db_url)

    if args.reset:
        print("\n⚠ Resetting document tables...")
        cur = conn.cursor()
        cur.execute("DELETE FROM document_chunks")
        cur.execute("DELETE FROM documents")
        conn.commit()
        cur.close()
        print("  done")

    ingest_to_db(manifest, conn, voyage_client, dry_run=False)
    conn.close()

    print("\n" + "=" * 60)
    print("✅ Ingestion complete")
    print("=" * 60)


if __name__ == "__main__":
    main()
